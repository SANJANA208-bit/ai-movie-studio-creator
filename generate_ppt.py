from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()
    
    # helper for slide layout
    # 0 = Title, 1 = Title and Content, 6 = Blank
    
    def add_slide(title, content_list=None, image_path=None, layout_idx=1):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # Set background to dark
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(10, 10, 20) # Very dark navy
        
        # Title styling
        title_shape = slide.shapes.title
        title_shape.text = title
        title_text_frame = title_shape.text_frame
        title_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_text_frame.paragraphs[0].font.bold = True
        title_text_frame.paragraphs[0].font.size = Pt(44)
        
        # Content styling
        if content_list:
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.word_wrap = True
            for item in content_list:
                p = tf.add_paragraph()
                p.text = item
                p.font.color.rgb = RGBColor(200, 200, 255) # Light blueish
                p.font.size = Pt(20)
                p.space_after = Pt(10)
        
        # Image placement
        if image_path and os.path.exists(image_path):
            # Place image on the right or bottom
            if content_list:
                # Shrink text box to make room for image on the right
                body_shape.width = Inches(4.5)
                left = Inches(5.0)
                top = Inches(1.5)
                width = Inches(4.5)
                # Height will scale
                slide.shapes.add_picture(image_path, left, top, width=width)
            else:
                # Full center image
                left = Inches(1.0)
                top = Inches(1.5)
                width = Inches(8.0)
                slide.shapes.add_picture(image_path, left, top, width=width)

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(15, 15, 30)
    
    title = slide.shapes.title
    title.text = "AI Movie Studio Creator"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Project-Based Assessment: Domain-Specific Generative AI Chatbot\nPresented by: [Your Name]"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 149, 237) # Cornflower Blue

    # 2. Project Overview & Objectives
    add_slide(
        "Project Overview & Objectives",
        [
            "Goal: Solve real-world creative bottlenecks using Generative AI.",
            "Domain: Film Concept Development & Screenwriting.",
            "Innovation: Automating the 'Blank Page' problem for storytellers.",
            "Societal Impact: Democratizing film-making tools for indie creators."
        ],
        "C:/Users/SANJANA/.gemini/antigravity/brain/c8482eff-e216-414b-9b7e-50ab9a4b5f06/media__1776242706269.png"
    )

    # 3. Domain Selection: Why Film?
    add_slide(
        "Domain Selection & Specificity",
        [
            "Focus: Tailored for screenwriters and producers.",
            "Originality: Domain-specific templates (Noir, Space Opera, Cyberpunk).",
            "Logic: Uses narrative structures (Hero's Journey, Conflict/Resolution).",
            "Constraints: Strictly filtered for movie-related queries only."
        ],
        "C:/Users/SANJANA/.gemini/antigravity/brain/c8482eff-e216-414b-9b7e-50ab9a4b5f06/media__1776242962807.png"
    )

    # 4. API & Security Configuration
    add_slide(
        "API & Security Configuration",
        [
            "Engine: Powered by Google Gemini GenAI APIs.",
            "Security: Managed via Environment Variables (.env).",
            "Performance: Low-latency responses with fallback mock data.",
            "Integration: Python Backend (FastAPI) + Google GenAI Client."
        ],
        "C:/Users/SANJANA/.gemini/antigravity/brain/c8482eff-e216-414b-9b7e-50ab9a4b5f06/api_config_screenshot_1776243246766.png"
    )

    # 5. Interactive Concept Design
    add_slide(
        "User-Centric Concept Design",
        [
            "Interactive Modals: Streamlined input for story ideas.",
            "Genre Control: Dynamic prompt building based on selection.",
            "Accessibility: Simple, intuitive UI for high-speed brainstorming."
        ],
        "C:/Users/SANJANA/.gemini/antigravity/brain/c8482eff-e216-414b-9b7e-50ab9a4b5f06/media__1776242770874.png"
    )

    # 6. AI Storytelling & Assistant
    add_slide(
        "AI-Driven Storytelling",
        [
            "Core Output: Title, Plot Summary, Characters, and Trailer Scripts.",
            "Story Assistant: Context-aware chat for real-time refinement.",
            "Transparency: Users can view the model's logic and constructed prompt."
        ],
        "C:/Users/SANJANA/.gemini/antigravity/brain/c8482eff-e216-414b-9b7e-50ab9a4b5f06/media__1776242808989.png"
    )

    # 7. Multi-modal Assets
    add_slide(
        "Multi-modal Concept Assets",
        [
            "Visuals: AI-generated character portraits and mood boards.",
            "Audio: Narrated trailer previews powered by ElevenLabs.",
            "Innovation: A one-stop studio for script and production prep."
        ],
        "C:/Users/SANJANA/.gemini/antigravity/brain/c8482eff-e216-414b-9b7e-50ab9a4b5f06/media__1776242782808.png"
    )

    # 8. Project Management
    add_slide(
        "Project Management & State",
        [
            "Consistency: Persistent project storage with metadata.",
            "Quick Start: Pre-defined industry-standard templates.",
            "Scalability: Easily add new genres and model instructions."
        ],
        "C:/Users/SANJANA/.gemini/antigravity/brain/c8482eff-e216-414b-9b7e-50ab9a4b5f06/media__1776242950742.png"
    )

    # 9. Deliverables & Export
    add_slide(
        "Deliverables & Output",
        [
            "Export: Full scripts downloadable as .txt for production.",
            "Prototypes: Fully functional web app (Vite + FastAPI).",
            "Documentation: Structured reports and presentation slides."
        ],
        "C:/Users/SANJANA/.gemini/antigravity/brain/c8482eff-e216-414b-9b7e-50ab9a4b5f06/media__1776242989254.png"
    )

    # 10. Evaluation Rubrics (30 Marks)
    add_slide(
        "Evaluation Rubrics",
        [
            "1. Problem ID & Innovation (10 Marks): Clarity, Originality, Impact.",
            "2. Technical Execution (10 Marks): API Design, Prompt Engineering.",
            "3. Presentation & Q&A (10 Marks): Demo effectiveness, Depth of insight.",
            "Total Score: 30 Marks"
        ]
    )

    prs.save("AI_Movie_Studio_Assessment.pptx")
    print("Presentation created successfully: AI_Movie_Studio_Assessment.pptx")

if __name__ == "__main__":
    create_presentation()
